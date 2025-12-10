"""Tool for exporting HTML slides to PPTX format using BeautifulSoup for HTML parsing."""

from __future__ import annotations

import json
import logging
import os
import re
from datetime import datetime
from typing import TYPE_CHECKING

from pydantic import Field

from sgr_deep_research.core.base_tool import BaseTool
from sgr_deep_research.core.models import AgentStatesEnum
from sgr_deep_research.core.models_presentation import PresentationContext

if TYPE_CHECKING:
    from sgr_deep_research.core.agent_definition import AgentConfig
    from sgr_deep_research.core.models import ResearchContext

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c * 2 for c in hex_color])
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def extract_color_from_html(html: str, css_property: str, default: str) -> str:
    """Extract color value from inline CSS in HTML."""
    # Look for the property in style attributes or CSS
    # Use word boundary or start to avoid matching "background-color" when looking for "color"
    if css_property == "color":
        # Match "color:" but not "background-color:" or other prefixed versions
        pattern = r'(?<![a-zA-Z-])color\s*:\s*([#\w]+)'
    else:
        pattern = rf'{css_property}\s*:\s*([#\w]+)'
    match = re.search(pattern, html, re.IGNORECASE)
    if match:
        return match.group(1)
    return default


class HTMLSlideParser:
    """Parse HTML slide content using BeautifulSoup."""

    def __init__(self, html_content: str):
        try:
            from bs4 import BeautifulSoup
            self.soup = BeautifulSoup(html_content, 'html.parser')
        except ImportError:
            # Fallback to basic regex parsing if BeautifulSoup not available
            self.soup = None
            self.html = html_content

    def get_title(self) -> str:
        """Extract slide title."""
        if self.soup:
            # Try to find title in slide-title class or h1
            title_elem = self.soup.find(class_='slide-title') or self.soup.find('h1')
            if title_elem:
                return title_elem.get_text(strip=True)
        else:
            # Regex fallback
            match = re.search(r'<h1[^>]*class="slide-title"[^>]*>(.*?)</h1>', self.html, re.DOTALL | re.IGNORECASE)
            if match:
                return re.sub(r'<[^>]+>', '', match.group(1)).strip()
        return "Untitled Slide"

    def get_content_elements(self) -> list[dict]:
        """Extract content elements from the slide."""
        elements = []

        if self.soup:
            # Find the slide-content div
            content_div = self.soup.find(class_='slide-content')
            if not content_div:
                content_div = self.soup.body or self.soup

            for elem in content_div.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'ul', 'ol', 'blockquote', 'div']):
                # Skip nested elements we'll process separately
                if elem.find_parent(['ul', 'ol']) and elem.name == 'li':
                    continue

                if elem.name in ['h1', 'h2', 'h3', 'h4']:
                    # Skip the main slide title
                    if 'slide-title' in elem.get('class', []):
                        continue
                    elements.append({
                        'type': 'heading',
                        'level': int(elem.name[1]),
                        'text': elem.get_text(strip=True)
                    })
                elif elem.name == 'p':
                    text = elem.get_text(strip=True)
                    if text:
                        elements.append({
                            'type': 'paragraph',
                            'text': text
                        })
                elif elem.name in ['ul', 'ol']:
                    items = []
                    for li in elem.find_all('li', recursive=False):
                        items.append(li.get_text(strip=True))
                    if items:
                        elements.append({
                            'type': 'list',
                            'items': items,
                            'ordered': elem.name == 'ol'
                        })
                elif elem.name == 'blockquote' or 'quote' in elem.get('class', []):
                    text = elem.get_text(strip=True)
                    if text:
                        elements.append({
                            'type': 'quote',
                            'text': text
                        })
                elif elem.name == 'div' and 'two-column' in elem.get('class', []):
                    # Handle two-column layout
                    columns = elem.find_all('div', recursive=False)
                    col_texts = [col.get_text(strip=True) for col in columns if col.get_text(strip=True)]
                    if col_texts:
                        elements.append({
                            'type': 'two_column',
                            'columns': col_texts
                        })
        else:
            # Basic regex fallback
            # Extract list items
            list_match = re.search(r'<ul[^>]*>(.*?)</ul>', self.html, re.DOTALL | re.IGNORECASE)
            if list_match:
                items = re.findall(r'<li[^>]*>(.*?)</li>', list_match.group(1), re.DOTALL | re.IGNORECASE)
                items = [re.sub(r'<[^>]+>', '', item).strip() for item in items]
                if items:
                    elements.append({'type': 'list', 'items': items, 'ordered': False})

            # Extract paragraphs
            for p_match in re.finditer(r'<p[^>]*>(.*?)</p>', self.html, re.DOTALL | re.IGNORECASE):
                text = re.sub(r'<[^>]+>', '', p_match.group(1)).strip()
                if text:
                    elements.append({'type': 'paragraph', 'text': text})

            # Extract headings
            for h_match in re.finditer(r'<h([2-4])[^>]*>(.*?)</h\1>', self.html, re.DOTALL | re.IGNORECASE):
                level = int(h_match.group(1))
                text = re.sub(r'<[^>]+>', '', h_match.group(2)).strip()
                if text:
                    elements.append({'type': 'heading', 'level': level, 'text': text})

        return elements


class ExportPresentationTool(BaseTool):
    """Export all created slides to a PPTX presentation file.

    This tool should be called as the FINAL step after all slides have been created
    using CreateSlideTool. It compiles all HTML slides into a single PowerPoint file.

    IMPORTANT: Only call this after you have created ALL necessary slides.
    """

    reasoning: str = Field(description="Why the presentation is ready to be exported")
    presentation_title: str = Field(description="Title for the presentation file")
    author: str = Field(default="SGR Presentation Agent", description="Author name for the presentation metadata")
    summary: str = Field(description="Brief summary of the presentation content and key takeaways")

    async def __call__(self, context: ResearchContext, config: AgentConfig, **_) -> str:
        # Check for required libraries
        try:
            from pptx import Presentation  # noqa: I001
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as e:
            error_msg = f"python-pptx is required for PPTX export. Install with: pip install python-pptx. Error: {e}"
            logger.error(error_msg)
            return json.dumps({"status": "error", "message": error_msg})

        # Get presentation context
        if context.custom_context is None or not isinstance(context.custom_context, PresentationContext):
            return json.dumps({
                "status": "error",
                "message": "No slides found. Create slides using CreateSlideTool first."
            })

        pres_context: PresentationContext = context.custom_context

        if not pres_context.slides:
            return json.dumps({
                "status": "error",
                "message": "No slides to export. Create at least one slide first."
            })

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Process each slide
        for slide_data in pres_context.slides:
            # Parse HTML content
            parser = HTMLSlideParser(slide_data.html_content)

            # Extract colors from HTML
            bg_color = extract_color_from_html(slide_data.html_content, 'background-color', '#ffffff')
            text_color = extract_color_from_html(slide_data.html_content, 'color', '#333333')
            accent_color = extract_color_from_html(slide_data.html_content, 'border.*color|accent', '#007bff')

            # Convert colors
            bg_rgb = hex_to_rgb(bg_color)
            text_rgb = hex_to_rgb(text_color)
            accent_rgb = hex_to_rgb(accent_color)

            # Add a blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Add background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_rgb)

            # Add title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data.title
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*text_rgb)

            # Add accent underline below title
            accent_line = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), Inches(1.35), Inches(2), Inches(0.06)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = RGBColor(*accent_rgb)
            accent_line.line.fill.background()

            # Get content elements
            elements = parser.get_content_elements()

            # Add content
            content_top = Inches(1.6)
            content_left = Inches(0.5)
            content_width = Inches(12.333)
            content_height = Inches(5.4)

            if elements:
                content_box = slide.shapes.add_textbox(
                    content_left, content_top, content_width, content_height
                )
                tf = content_box.text_frame
                tf.word_wrap = True

                first_para = True
                for elem in elements:
                    elem_type = elem.get('type')

                    if elem_type == 'heading':
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        p.text = elem['text']
                        level = elem.get('level', 2)
                        p.font.size = Pt(28 - (level - 2) * 4)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(*text_rgb)
                        p.space_after = Pt(12)

                    elif elem_type == 'list':
                        for item in elem.get('items', []):
                            if first_para:
                                p = tf.paragraphs[0]
                                first_para = False
                            else:
                                p = tf.add_paragraph()
                            bullet = "• " if not elem.get('ordered') else f"{elem['items'].index(item) + 1}. "
                            p.text = f"{bullet}{item}"
                            p.font.size = Pt(20)
                            p.font.color.rgb = RGBColor(*text_rgb)
                            p.space_after = Pt(8)
                            p.level = 0

                    elif elem_type == 'paragraph':
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        p.text = elem['text']
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(*text_rgb)
                        p.space_after = Pt(10)

                    elif elem_type == 'quote':
                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()
                        p.text = f'"{elem["text"]}"'
                        p.font.size = Pt(22)
                        p.font.italic = True
                        p.font.color.rgb = RGBColor(*accent_rgb)
                        p.space_after = Pt(15)

                    elif elem_type == 'two_column':
                        # For two-column, we'll just add them as separate paragraphs
                        for col_text in elem.get('columns', []):
                            if first_para:
                                p = tf.paragraphs[0]
                                first_para = False
                            else:
                                p = tf.add_paragraph()
                            p.text = col_text
                            p.font.size = Pt(18)
                            p.font.color.rgb = RGBColor(*text_rgb)
                            p.space_after = Pt(10)

            # Add speaker notes if present
            if slide_data.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data.speaker_notes

            # Add slide number
            slide_num_box = slide.shapes.add_textbox(
                Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.35)
            )
            slide_num_frame = slide_num_box.text_frame
            slide_num_para = slide_num_frame.paragraphs[0]
            slide_num_para.text = str(slide_data.slide_number)
            slide_num_para.font.size = Pt(12)
            slide_num_para.font.color.rgb = RGBColor(*accent_rgb)
            slide_num_para.alignment = PP_ALIGN.RIGHT

        # Save the presentation
        reports_dir = config.execution.reports_dir
        os.makedirs(reports_dir, exist_ok=True)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in self.presentation_title if c.isalnum() or c in (" ", "-", "_"))[:50]
        pptx_filename = f"{timestamp}_{safe_title}.pptx"
        pptx_filepath = os.path.join(reports_dir, pptx_filename)

        # Save PPTX
        prs.save(pptx_filepath)

        # Also save HTML files for reference
        html_dir = os.path.join(reports_dir, f"{timestamp}_{safe_title}_html")
        os.makedirs(html_dir, exist_ok=True)

        for slide_data in pres_context.slides:
            html_filename = f"slide_{slide_data.slide_number:02d}.html"
            html_filepath = os.path.join(html_dir, html_filename)
            with open(html_filepath, "w", encoding="utf-8") as f:
                f.write(slide_data.html_content)

        # Update agent state
        context.state = AgentStatesEnum.COMPLETED
        context.execution_result = f"Presentation exported to {pptx_filepath}"

        result = {
            "status": "success",
            "title": self.presentation_title,
            "total_slides": len(pres_context.slides),
            "pptx_filepath": pptx_filepath,
            "html_folder": html_dir,
            "summary": self.summary,
            "slides": [{"number": s.slide_number, "title": s.title} for s in pres_context.slides],
        }

        logger.info(
            f"📊 PRESENTATION EXPORTED:\n"
            f"   📝 Title: '{self.presentation_title}'\n"
            f"   📑 Total slides: {len(pres_context.slides)}\n"
            f"   💾 PPTX: {pptx_filepath}\n"
            f"   🌐 HTML: {html_dir}\n"
            f"   ✅ Summary: {self.summary[:100]}...\n"
        )

        return json.dumps(result, indent=2, ensure_ascii=False)

